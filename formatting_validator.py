"""
Formatting validator for JoVE PPT V6.
Checks formatting rules that can be verified from the PPTX object model.
"""

import re
from typing import Dict, List
from pptx import Presentation
from pptx.util import Inches

FOOTER_TEXT = "Copyright © 2026 MyJoVE Corporation. All rights reserved"
BRAND_BLUE_HEX = "6D9EEB"
WHITE_HEX = "FFFFFF"
FORBIDDEN_RE = re.compile(r"\b(writer|author|reviewer|prepared\s*by|created\s*by)\s*[:\-]", re.I)
MARKDOWN_RE = re.compile(r"(\*\*|__|`|\[INSERT IMAGE\]|TODO|PLACEHOLDER)", re.I)


def _rgb_hex(rgb) -> str:
    try:
        return str(rgb).upper()
    except Exception:
        return ""


def _shape_text(shape) -> str:
    try:
        if hasattr(shape, "text"):
            return shape.text or ""
    except Exception:
        return ""
    return ""


def _run_fonts(shape) -> List[Dict]:
    out = []
    if not getattr(shape, "has_text_frame", False):
        return out
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            size = None
            try:
                size = run.font.size.pt if run.font.size else None
            except Exception:
                pass
            name = None
            try:
                name = run.font.name
            except Exception:
                pass
            out.append({"text": run.text, "size": size, "font": name})
    return out


def validate_pptx_formatting(pptx_path: str) -> Dict:
    prs = Presentation(pptx_path)
    findings = []
    checks = 0
    passes = 0

    def check(condition: bool, slide_num: int, rule: str, detail: str = ""):
        nonlocal checks, passes
        checks += 1
        if condition:
            passes += 1
        else:
            findings.append({"slide": slide_num, "rule": rule, "detail": detail})

    for idx, slide in enumerate(prs.slides, start=1):
        texts = "\n".join(_shape_text(s) for s in slide.shapes)
        check(not FORBIDDEN_RE.search(texts), idx, "No writer/author/reviewer metadata", "Forbidden metadata text found")
        check(not MARKDOWN_RE.search(texts), idx, "No markdown/placeholders", "Markdown or placeholder text found")
        check(FOOTER_TEXT in texts, idx, "Footer present", "Exact copyright footer missing")
        check(str(idx) in texts, idx, "Slide number present", "Expected slide number missing")

        logo_candidates = [
            s for s in slide.shapes
            if getattr(s, "shape_type", None) == 13
            and s.left >= Inches(17.65)
            and s.top <= Inches(0.35)
        ]
        check(len(logo_candidates) >= 1, idx, "Slide-level JoVE logo present at fixed position", "Missing fixed top-right logo")
        if logo_candidates:
            logo = logo_candidates[0]
            check(abs((logo.left / Inches(1)) - 17.8) <= 0.12, idx, "Logo fixed x-position", f"x={logo.left/Inches(1):.2f}")
            check(abs((logo.top / Inches(1)) - 0.15) <= 0.12, idx, "Logo fixed y-position", f"y={logo.top/Inches(1):.2f}")

        all_runs = []
        for shape in slide.shapes:
            all_runs.extend(_run_fonts(shape))
        font_names = [r.get("font") for r in all_runs if r.get("text") and r.get("font")]
        if font_names:
            check(all("Roboto" in str(f) or "Helvetica" in str(f) or "Arial" in str(f) for f in font_names), idx, "Allowed fonts only", ", ".join(sorted(set(font_names)))[:100])

        # Title size check: largest non-footer text on cover should be about 102; content slides should be 48.
        large_runs = [r for r in all_runs if r.get("text") and r.get("size") and FOOTER_TEXT not in r.get("text", "")]
        if large_runs:
            max_size = max(r["size"] for r in large_runs if r["size"])
            if idx == 1:
                check(max_size >= 96, idx, "Cover title size near 102pt", f"Max font size: {max_size}")
            else:
                # Some answer subheadings are 32, but a slide should still have a 48pt heading.
                check(max_size >= 46, idx, "Slide heading size near 48pt", f"Max font size: {max_size}")

        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            table_bottom = shape.top + shape.height
            check(shape.top >= Inches(1.95), idx, "Table starts below title safe zone", f"top={shape.top/Inches(1):.2f}in")
            check(table_bottom <= Inches(10.25), idx, "Table stays above footer zone", f"bottom={table_bottom/Inches(1):.2f}in")
            tbl = shape.table
            # Header row must be blue.
            for c in range(len(tbl.columns)):
                try:
                    fill = _rgb_hex(tbl.cell(0, c).fill.fore_color.rgb)
                    check(fill == BRAND_BLUE_HEX, idx, "All table header cells blue", f"header cell {c+1} fill={fill}")
                except Exception:
                    findings.append({"slide": idx, "rule": "All table header cells blue", "detail": f"Could not inspect header cell {c+1}"})
            # Approved override: body rows all white.
            for r in range(1, len(tbl.rows)):
                for c in range(len(tbl.columns)):
                    try:
                        fill = _rgb_hex(tbl.cell(r, c).fill.fore_color.rgb)
                        check(fill == WHITE_HEX, idx, "All table body cells white", f"row {r+1}, col {c+1} fill={fill}")
                    except Exception:
                        pass

        # Right image position check for non-table picture-heavy slides.
        pic_shapes = [s for s in slide.shapes if getattr(s, "shape_type", None) == 13]
        # Ignore the top-right logo by filtering tiny top images.
        content_pics = [s for s in pic_shapes if not (s.left >= Inches(17.4) and s.top <= Inches(0.8))]
        if idx == 1 and len(content_pics) > 1:
            boxes = [(round(s.left/Inches(1), 2), round(s.top/Inches(1), 2), round(s.width/Inches(1), 2), round(s.height/Inches(1), 2)) for s in content_pics]
            check(len(set(boxes)) == len(boxes), idx, "Cover does not duplicate identical image boxes", str(boxes))
        if content_pics and not any(getattr(s, "has_table", False) for s in slide.shapes):
            # Most content pictures should start in right column or cover stack.
            for pic in content_pics:
                if idx == 1:
                    check(pic.left >= Inches(10.8), idx, "Cover images in right column", f"picture x={pic.left/Inches(1):.2f}")
                else:
                    check(pic.left >= Inches(10.8), idx, "Right-column image x-position", f"picture x={pic.left/Inches(1):.2f}")
                    check(pic.left + pic.width <= Inches(19.25), idx, "Image within right margin", f"right={(pic.left+pic.width)/Inches(1):.2f}")

    score = round((passes / checks) * 100, 2) if checks else 0.0
    return {
        "formatting_score": score,
        "checks": checks,
        "passes": passes,
        "failures": len(findings),
        "findings": findings[:200],
        "target_met": score >= 95.0,
    }
