"""
JoVE PPT Builder V6 - strict formatting templates.

AI provides content only. This builder owns formatting decisions.
The goal is deterministic adherence to the JoVE presentation guide and the
approved project overrides.
"""

import os
import re
import tempfile
from typing import Iterable, List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

try:
    from PIL import Image, ImageDraw
except Exception:  # Pillow should be installed by requirements.
    Image = None

try:
    from style_guide import (
        FOOTER_TEXT, FONT_PRIMARY,
        PRIMARY_DARK, BRAND_BLUE, BRAND_BLUE_LIGHT, WHITE, MID_GRAY,
        BLACK, DARK_GRAY, LIGHT_GRAY, SLIDE_W, SLIDE_H, MARGIN, RIGHT_X, RIGHT_W
    )
except Exception:
    FOOTER_TEXT = "Copyright © 2026 MyJoVE Corporation. All rights reserved"
    FONT_PRIMARY = "Roboto"
    PRIMARY_DARK = RGBColor(0x24, 0x29, 0x2F)
    BRAND_BLUE = RGBColor(0x6D, 0x9E, 0xEB)
    BRAND_BLUE_LIGHT = RGBColor(0xA4, 0xC2, 0xF4)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    MID_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    DARK_GRAY = RGBColor(0x4B, 0x55, 0x63)
    LIGHT_GRAY = RGBColor(0x85, 0x85, 0x85)
    SLIDE_W = Inches(20)
    SLIDE_H = Inches(11.25)
    MARGIN = Inches(0.75)
    RIGHT_X = Inches(11.0)
    RIGHT_W = Inches(8.25)

C_TEXT_DARK = PRIMARY_DARK
C_SUBTITLE = DARK_GRAY
C_COPYRIGHT = BLACK
C_ACCENT_BLUE = BRAND_BLUE
C_TABLE_HEADER = BRAND_BLUE
C_TABLE_ROW = WHITE  # approved override: all table body rows are white
C_WHITE = WHITE
C_BORDER = MID_GRAY
C_LIGHT_PANEL = RGBColor(0xF3, 0xF6, 0xFB)

FONT = FONT_PRIMARY

# Guide-driven layout constants.
SLIDE_SAFE_TOP = Inches(0.35)
SLIDE_SAFE_BOTTOM = Inches(10.35)
TITLE_SAFE_TOP_TABLE = Inches(2.05)
TITLE_BOX_H = Inches(1.25)
COVER_TITLE_MAX_W = Inches(10.9)
COVER_SINGLE_IMAGE_H = Inches(5.2)
LEFT = Inches(0.75)
TEXT_W = Inches(9.65)
GUTTER = Inches(0.30)
IMG_L = RIGHT_X
IMG_T = Inches(1.55)
IMG_W = RIGHT_W
IMG_H = Inches(8.45)
LOGO_L = Inches(18.444)
LOGO_T = Inches(0.326)
LOGO_W = Inches(1.087)
LOGO_H = Inches(0.551)
FOOTER_T = Inches(10.64)
FOOTER_H = Inches(0.28)

# Required font sizes.
FS_COVER_TITLE = 102
FS_SLIDE_TITLE = 48
FS_BODY = 30
FS_BODY_SECONDARY = 24
FS_TABLE_HEADER = 28
FS_TABLE_BODY = 26
FS_DISCUSSION_BADGE = 20
FS_CAPTION = 14
FS_FOOTER = 11

FORBIDDEN_LINE_RE = re.compile(r"^\s*(writer|author|reviewer|prepared\s*by|created\s*by|presenter|date|file\s*name|source\s*file|pagetext|page\s*text|script|transcript|transcription)\s*[:\-]", re.I)

COMMON_TEXT_FIXES = {
    "Funtion": "Function",
    "funtion": "function",
}


def _apply_common_text_fixes(value: str) -> str:
    for bad, good in COMMON_TEXT_FIXES.items():
        value = value.replace(bad, good)
    return value


def _clean_text(text) -> str:
    """Remove markdown and forbidden metadata before writing to PPT."""
    if text is None:
        return ""
    value = str(text)
    value = value.replace("**", "")
    value = value.replace("__", "")
    value = value.replace("`", "")
    value = re.sub(r"\[(INSERT IMAGE|TODO|PLACEHOLDER|IMAGE)\]", "", value, flags=re.I)
    lines = []
    for line in value.splitlines():
        if FORBIDDEN_LINE_RE.search(line):
            continue
        lines.append(line.strip())
    value = "\n".join(line for line in lines if line)
    value = _apply_common_text_fixes(value)
    return re.sub(r"[ \t]+", " ", value).strip()



def _safe_slide_title(title: str, body_text: str = "") -> str:
    """Prevent numeric lesson IDs or metadata from becoming visible slide titles."""
    raw = _clean_text(title or "")
    if not raw or re.fullmatch(r"(lesson\s*)?\d{4,8}", raw.strip(), flags=re.I):
        body = _clean_text(body_text or "")
        m = re.match(r"([A-Z][A-Za-z\- ]{2,35}?)(?:\s+are|\s+is|\s+include|\s+consist|\s+form|\s+have)\b", body)
        if m:
            return m.group(1).strip()
        words = re.findall(r"[A-Za-z][A-Za-z\-]+", body)[:4]
        return " ".join(words).title() if words else "Core Concept"
    return raw


def _fit_title_font(title: str, target=FS_SLIDE_TITLE, min_size=34) -> int:
    title = str(title or "")
    if len(title) <= 58:
        return target
    if len(title) <= 72:
        return max(min_size, target - 6)
    if len(title) <= 88:
        return max(min_size, target - 10)
    return max(min_size, target - 14)


def _fit_cover_font(title: str) -> int:
    title = str(title or "")
    if len(title) <= 14:
        return FS_COVER_TITLE
    if len(title) <= 22:
        return 94
    if len(title) <= 32:
        return 84
    return 76


def _shorten_words(text: str, max_words: int) -> str:
    """Shorten for compact cells/labels while keeping a clean ending."""
    text = _clean_text(text)
    words = text.split()
    if len(words) <= max_words:
        return text
    return " ".join(words[:max_words]).rstrip(".,;:") + "."


def _sentence_list(text: str):
    text = _clean_text(text)
    if not text:
        return []
    # Keep question marks/exclamation marks with the sentence.
    parts = re.split(r"(?<=[.!?])\s+", text)
    return [p.strip() for p in parts if p.strip()]


def _is_complete_sentence(text: str) -> bool:
    text = _clean_text(text)
    return bool(text and re.search(r"[.!?]$", text) and not re.search(r"\b(and|or|for|with|to|of|the|their|a|an|relative|growing|associated)\.$", text, re.I))


def _complete_sentence_text(primary: str, fallback: str = "", max_words: int = 32, prefer_question: bool = False) -> str:
    """Return visible text that ends naturally.

    This prevents truncated discussion lines such as 'What evidence supports.'
    or 'Answer: ... for.' from appearing on slides.
    """
    candidates = []
    for source in (primary, fallback):
        source = _clean_text(source)
        if not source:
            continue
        # If it is already short and complete, use it.
        if len(source.split()) <= max_words and _is_complete_sentence(source):
            candidates.append(source)
        candidates.extend(_sentence_list(source))
    for cand in candidates:
        if len(cand.split()) <= max_words and _is_complete_sentence(cand):
            return cand
    # Build from complete sentences until budget is reached.
    assembled = []
    count = 0
    for cand in candidates:
        if not _is_complete_sentence(cand):
            continue
        wc = len(cand.split())
        if count + wc > max_words and assembled:
            break
        if wc <= max_words or not assembled:
            assembled.append(cand)
            count += wc
    if assembled:
        return " ".join(assembled)
    # Last resort: shorten without creating a false sentence if source is a question.
    base = _clean_text(primary or fallback)
    words = base.split()[:max_words]
    text = " ".join(words).rstrip(".,;:")
    if not text:
        return ""
    return text + ("?" if prefer_question else ".")


def _font(run, size_pt, bold=False, italic=False, color=None):
    """Apply the same JoVE font family everywhere.

    Headings, body text, table text, glossary text, discussion text, footer,
    and slide numbers all use the same FONT value. They may differ by size
    and weight, but never by font family.
    """
    run.font.name = FONT
    try:
        rFonts = run._r.rPr.rFonts
        for attr in ("ascii", "hAnsi", "eastAsia", "cs"):
            rFonts.set(f"{{http://schemas.openxmlformats.org/wordprocessingml/2006/main}}{attr}", FONT)
    except Exception:
        pass
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color


def _lock_text_frame(tf, vertical_anchor=MSO_ANCHOR.TOP, wrap=True):
    """Normalize all text frames so Google Slides/PowerPoint do not
    auto-fit, justify, or distribute text in a way that creates stretched
    word spacing. This is intentionally global, not discussion-only.
    """
    tf.word_wrap = wrap
    tf.vertical_anchor = vertical_anchor
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        tf.margin_left = Inches(0.0)
        tf.margin_right = Inches(0.0)
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
    except Exception:
        pass


def _format_paragraph(p, align=PP_ALIGN.LEFT, space_before=0, space_after=0, line_spacing=1.0):
    """Force normal paragraph behavior everywhere.

    We never use JUSTIFY, DISTRIBUTE, or Thai-distributed alignment because
    those can create the stretched-spacing issue seen in the discussion slide.
    Tables may be centered intentionally, but still use normal line spacing.
    """
    p.alignment = align
    try:
        p.space_before = Pt(space_before)
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
    except Exception:
        pass


def _tb(slide, left, top, width, height, text, size_pt,
        bold=False, italic=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    _lock_text_frame(tf, vertical_anchor=MSO_ANCHOR.TOP, wrap=wrap)
    p = tf.paragraphs[0]
    _format_paragraph(p, align=align, space_before=0, space_after=0, line_spacing=1.0)
    run = p.add_run()
    run.text = _clean_text(text)
    _font(run, size_pt, bold, italic, color or C_TEXT_DARK)
    return box


def _white_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_WHITE


def _logo(slide, logo_path):
    """Place slide-level JoVE logo in one exact location on every slide."""
    if logo_path and os.path.exists(logo_path):
        try:
            pic = slide.shapes.add_picture(logo_path, LOGO_L, LOGO_T, width=LOGO_W, height=LOGO_H)
            return pic
        except Exception:
            # Do not create inconsistent fallback logos.
            return None



def _copyright(slide):
    _tb(slide, Inches(5.35), FOOTER_T, Inches(9.3), FOOTER_H,
        FOOTER_TEXT, FS_FOOTER, color=C_COPYRIGHT, align=PP_ALIGN.CENTER)


def _slide_number(slide, number):
    if number is None:
        return
    _tb(slide, Inches(18.65), Inches(10.62), Inches(0.55), Inches(0.28),
        str(number), FS_FOOTER, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def _notes(slide, text):
    text = _clean_text(text)
    if text:
        slide.notes_slide.notes_text_frame.text = text


def _base_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _prepare_presentation_image(image_path: str) -> str:
    """Light local cleanup so selected frames are presentation-safe.

    This is the local, no-cost visual pass. Full Natural Selection-style
    transformation is handled upstream when optional JOVE_AI_VISUALS=1 is enabled
    in the pipeline. This function still improves contrast/color/sharpness and
    removes obvious uniform borders while preserving the source frame.
    """
    if Image is None or not image_path or not os.path.exists(image_path):
        return image_path
    try:
        from PIL import ImageOps, ImageEnhance, ImageChops
        cache_dir = os.path.join(tempfile.gettempdir(), "jove_presentation_images")
        os.makedirs(cache_dir, exist_ok=True)
        base = os.path.splitext(os.path.basename(image_path))[0]
        out_path = os.path.join(cache_dir, f"{base}_presentation.png")
        if os.path.exists(out_path) and os.path.getmtime(out_path) >= os.path.getmtime(image_path):
            return out_path
        with Image.open(image_path) as img:
            img = img.convert("RGB")
            # Crop near-uniform border/background if present.
            bg = Image.new(img.mode, img.size, img.getpixel((0, 0)))
            diff = ImageChops.difference(img, bg)
            bbox = diff.getbbox()
            if bbox:
                l, t, r, b = bbox
                # Only crop if it does not remove too much content.
                if (r - l) > img.width * 0.55 and (b - t) > img.height * 0.55:
                    pad_x = int(img.width * 0.015)
                    pad_y = int(img.height * 0.015)
                    img = img.crop((max(0, l - pad_x), max(0, t - pad_y), min(img.width, r + pad_x), min(img.height, b + pad_y)))
            img = ImageOps.autocontrast(img, cutoff=1)
            img = ImageEnhance.Color(img).enhance(1.08)
            img = ImageEnhance.Sharpness(img).enhance(1.12)
            img.save(out_path, "PNG", optimize=True)
        return out_path
    except Exception:
        return image_path


def _add_image_contain(slide, image_path, left, top, width, height):
    if not image_path or not os.path.exists(image_path):
        raise ValueError("A valid local image_path is required for every image-bearing slide. No placeholders are allowed.")
    image_path = _prepare_presentation_image(image_path)
    if Image is None:
        slide.shapes.add_picture(image_path, left, top, width=width)
        return
    try:
        with Image.open(image_path) as img:
            iw, ih = img.size
        if iw <= 0 or ih <= 0:
            slide.shapes.add_picture(image_path, left, top, width=width)
            return
        box_ratio = float(width) / float(height)
        img_ratio = float(iw) / float(ih)
        if img_ratio >= box_ratio:
            final_w = width
            final_h = int(width / img_ratio)
        else:
            final_h = height
            final_w = int(height * img_ratio)
        x = left + int((width - final_w) / 2)
        y = top + int((height - final_h) / 2)
        slide.shapes.add_picture(image_path, x, y, width=final_w, height=final_h)
    except Exception:
        slide.shapes.add_picture(image_path, left, top, width=width)



def _discussion_icon_path():
    """Create a deterministic outline speech-bubble icon as a tiny transparent PNG."""
    if Image is None:
        return None
    try:
        path = os.path.join(tempfile.gettempdir(), "jove_discussion_icon_outline.png")
        if os.path.exists(path):
            return path

        img = Image.new("RGBA", (96, 96), (255, 255, 255, 0))
        draw = ImageDraw.Draw(img)
        blue = (0x6D, 0x9E, 0xEB, 255)
        # Rounded rectangle bubble.
        draw.rounded_rectangle((14, 18, 78, 66), radius=9, outline=blue, width=7)
        # Tail.
        draw.line((32, 66, 22, 82, 46, 66), fill=blue, width=7, joint="curve")
        img.save(path)
        return path
    except Exception:
        return None


def _short_caption(text: str, max_words: int = 5) -> str:
    words = re.findall(r"[A-Za-z0-9][A-Za-z0-9\-]*", _clean_text(text or ""))[:max_words]
    return " ".join(words).strip()


def _figure_legend(slide, text: str):
    caption = _short_caption(text, max_words=5)
    if caption:
        _tb(slide, IMG_L, Inches(9.82), IMG_W, Inches(0.34),
            caption, FS_CAPTION, italic=True, color=DARK_GRAY, align=PP_ALIGN.CENTER, wrap=False)


def _transition_caption(slide, text: str):
    caption = _complete_sentence_text(text, "", max_words=20) if text else ""
    if caption:
        _tb(slide, LEFT, Inches(9.76), TEXT_W, Inches(0.55),
            caption, FS_CAPTION, italic=True, color=DARK_GRAY, align=PP_ALIGN.LEFT, wrap=True)


def _image(slide, image_path, figure_legend=None):
    if figure_legend:
        _add_image_contain(slide, image_path, IMG_L, IMG_T, IMG_W, Inches(8.05))
        _figure_legend(slide, figure_legend)
    else:
        _add_image_contain(slide, image_path, IMG_L, IMG_T, IMG_W, IMG_H)
    return True


def _body_text(slide, body_text, top, max_words=68):
    body_text = _shorten_words(body_text, max_words)
    box = slide.shapes.add_textbox(LEFT, top, TEXT_W, SLIDE_SAFE_BOTTOM - top)
    tf = box.text_frame
    tf.clear()
    _lock_text_frame(tf, vertical_anchor=MSO_ANCHOR.TOP, wrap=True)

    lines = [ln.strip() for ln in body_text.split("\n") if ln.strip()] or [""]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        _format_paragraph(p, align=PP_ALIGN.LEFT, space_before=(10 if idx > 0 else 0), space_after=0, line_spacing=1.0)
        segments = re.split(r'\*\*(.+?)\*\*', line)
        for i, seg in enumerate(segments):
            if not seg:
                continue
            run = p.add_run()
            run.text = _clean_text(seg)
            _font(run, FS_BODY, bold=(i % 2 == 1), color=C_TEXT_DARK)


def _normalize_table(headers, rows, table_kind=None):
    headers = [_clean_text(h) for h in (headers or [])]
    rows = rows or []
    if not headers:
        headers = ["Concept", "Definition/Meaning", "Example/Application"]

    n_cols = max(1, len(headers))
    normalized = []
    for row in rows:
        row = [_clean_text(x) for x in list(row)]
        if len(row) < n_cols:
            row += [""] * (n_cols - len(row))
        normalized.append(row[:n_cols])
    if not normalized:
        normalized = [[""] * n_cols]
    return headers, normalized



def _set_cell_border_blue(cell):
    """Apply native PowerPoint blue borders to every side of an actual table cell."""
    try:
        tcPr = cell._tc.get_or_add_tcPr()
        for edge in ("lnL", "lnR", "lnT", "lnB"):
            tag = qn(f"a:{edge}")
            ln = tcPr.find(tag)
            if ln is None:
                ln = OxmlElement(f"a:{edge}")
                tcPr.append(ln)

            # Reset existing hidden/no-line settings.
            for child in list(ln):
                ln.remove(child)

            ln.set("w", "19050")  # 1.5 pt
            ln.set("cap", "flat")
            ln.set("cmpd", "sng")
            ln.set("algn", "ctr")

            solid = OxmlElement("a:solidFill")
            srgb = OxmlElement("a:srgbClr")
            srgb.set("val", "6D9EEB")
            solid.append(srgb)
            ln.append(solid)

            dash = OxmlElement("a:prstDash")
            dash.set("val", "solid")
            ln.append(dash)

            round_join = OxmlElement("a:round")
            ln.append(round_join)
    except Exception:
        pass


def _cell_text(cell, text, size, bold=False, align=PP_ALIGN.CENTER, color=None):
    try:
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    tf = cell.text_frame
    tf.clear()
    _lock_text_frame(tf, vertical_anchor=MSO_ANCHOR.MIDDLE, wrap=True)
    p = tf.paragraphs[0]
    _format_paragraph(p, align=align, space_before=0, space_after=0, line_spacing=1.0)
    run = p.add_run()
    run.text = _clean_text(text)
    _font(run, size, bold=bold, color=color or C_TEXT_DARK)




def _estimate_lines_for_box(text: str, width_emu, font_size: int) -> int:
    """Conservative text-line estimate to keep text inside visual table cells."""
    text = _clean_text(text)
    if not text:
        return 1
    width_in = max(0.6, float(width_emu) / 914400.0)
    chars_per_line = max(8, int(width_in * 150 / max(10, font_size)))
    lines = 0
    for part in text.splitlines() or [text]:
        words = part.split()
        current = 0
        for word in words:
            add = len(word) + (1 if current else 0)
            if current + add > chars_per_line:
                lines += 1
                current = len(word)
            else:
                current += add
        lines += 1 if current or not words else 0
    return max(1, lines)


def _fit_shape_font_size(text: str, width_emu, height_emu, target_size: int, min_size: int = 14) -> int:
    """Reduce font size only when needed so text stays inside the cell."""
    height_pt = (float(height_emu) / 914400.0) * 72.0
    for size in range(int(target_size), int(min_size) - 1, -1):
        lines = _estimate_lines_for_box(text, width_emu, size)
        needed = lines * size * 1.12
        if needed <= height_pt * 0.86:
            return size
    return min_size


def _shape_cell_text(shape, text, size, bold=False, align=PP_ALIGN.CENTER, color=None, min_size=14):
    """Write centered text into a visual table cell without overflowing."""
    tf = shape.text_frame
    tf.clear()
    _lock_text_frame(tf, vertical_anchor=MSO_ANCHOR.MIDDLE, wrap=True)
    try:
        tf.margin_left = Inches(0.07)
        tf.margin_right = Inches(0.07)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
    except Exception:
        pass

    fitted_size = _fit_shape_font_size(text, shape.width, shape.height, size, min_size=min_size)
    p = tf.paragraphs[0]
    _format_paragraph(p, align=align, space_before=0, space_after=0, line_spacing=0.95)
    run = p.add_run()
    run.text = _clean_text(text)
    _font(run, fitted_size, bold=bold, color=color or C_TEXT_DARK)


def _table_cell_limit(ci: int, n_cols: int, has_images: bool) -> int:
    # Keep useful explanation while preventing crowded cells.
    if ci == 0:
        return 10
    if has_images and ci == n_cols - 1:
        return 0
    if n_cols >= 4:
        return 28
    return 34



def _drop_fully_empty_text_columns(headers, rows):
    """Remove columns whose body cells are all blank. Used only for text-only summary tables."""
    if not headers:
        return headers, rows
    keep = []
    for ci, header in enumerate(headers):
        header_text = str(header or "").strip()
        has_any_body = any(ci < len(row) and str(row[ci] or "").strip() for row in rows)
        if header_text and has_any_body:
            keep.append(ci)
    if not keep:
        return headers, rows
    return [headers[i] for i in keep], [[row[i] if i < len(row) else "" for i in keep] for row in rows]


def _add_table(slide, headers, rows, left, top, width, max_height=Inches(7.95), row_image_paths=None, max_rows=4):
    """Build a visually explicit table as cell rectangles.

    This replaces the unreliable native table-border rendering. Each visible
    cell is its own PowerPoint rectangle with JoVE-blue outline, so the output
    cannot appear as one big borderless block.
    """
    headers, rows = _normalize_table(headers, rows)
    row_image_paths = list(row_image_paths or [])
    if not row_image_paths:
        headers, rows = _drop_fully_empty_text_columns(headers, rows)

    if len(rows) > max_rows:
        # Tables are intentionally kept on one slide. Upstream prompt/cleanup must keep
        # normal tables readable (usually 3-4 rows) instead of splitting by default.
        rows = rows[:max_rows]
        row_image_paths = row_image_paths[:max_rows]

    # Image column is added only when EVERY row has a valid unique image.
    # This prevents blank Image cells and avoids repeated fallback images.
    valid_row_images = [p for p in row_image_paths[:len(rows)] if p and os.path.exists(p)]
    unique_row_images = {os.path.abspath(p) for p in valid_row_images}
    has_row_images = bool(rows) and len(valid_row_images) == len(rows) and len(unique_row_images) == len(rows)
    if not has_row_images:
        row_image_paths = []
    lower_headers = [str(h).strip().lower() for h in headers]
    has_image_col = any(h in {"image", "visual", "figure"} for h in lower_headers)
    rows = [list(row) for row in rows]
    if has_row_images and not has_image_col:
        headers.append("Image")
        rows = [row + [""] for row in rows]
    elif has_image_col:
        if not has_row_images:
            image_idx = next((i for i, h in enumerate(lower_headers) if h in {"image", "visual", "figure"}), len(headers) - 1)
            headers = [h for i, h in enumerate(headers) if i != image_idx]
            rows = [[v for i, v in enumerate(row) if i != image_idx] for row in rows]
        else:
            for row in rows:
                if len(row) < len(headers):
                    row.extend([""] * (len(headers) - len(row)))
                row[-1] = ""

    n_rows = len(rows) + 1
    n_cols = len(headers)
    max_allowed_height = max(0, SLIDE_SAFE_BOTTOM - top)
    height = min(max_height, max_allowed_height)
    row_h = height / max(1, n_rows)

    if has_row_images:
        image_col_w = int(width * 0.18)
        remaining = int(width) - image_col_w
        if n_cols == 4:
            col_widths = [int(remaining * 0.23), int(remaining * 0.42), int(remaining * 0.35), image_col_w]
        elif n_cols == 5:
            col_widths = [int(remaining * 0.18), int(remaining * 0.28), int(remaining * 0.29), int(remaining * 0.25), image_col_w]
        else:
            col_widths = [int(remaining / max(1, n_cols - 1))] * (n_cols - 1) + [image_col_w]
        drift = int(width) - sum(col_widths)
        if len(col_widths) >= 2:
            col_widths[-2] += drift
    else:
        # Summary tables have no Image column; all existing columns are text columns.
        col_widths = [int(width / max(1, n_cols))] * n_cols
        drift = int(width) - sum(col_widths)
        if col_widths:
            col_widths[-1] += drift

    header_size = FS_TABLE_HEADER if n_cols <= 3 else 24
    body_size = FS_TABLE_BODY if n_cols <= 3 else 24

    # Header cells.
    y = top
    x = left
    for ci, header in enumerate(headers):
        w = col_widths[ci]
        cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, row_h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TABLE_HEADER
        cell.line.color.rgb = C_TABLE_HEADER
        cell.line.width = Pt(1.5)
        _shape_cell_text(cell, _shorten_words(header, 8), header_size, bold=True, align=PP_ALIGN.CENTER, color=C_WHITE)
        x += w

    # Body cells with visible JoVE-blue borders on every cell.
    for ri, row in enumerate(rows):
        y = top + row_h * (ri + 1)
        x = left
        for ci in range(n_cols):
            w = col_widths[ci]
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, row_h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE
            cell.line.color.rgb = C_TABLE_HEADER
            cell.line.width = Pt(1.5)

            is_image_col = has_row_images and ci == n_cols - 1
            if not is_image_col:
                val = row[ci] if ci < len(row) else ""
                limit = _table_cell_limit(ci, n_cols, has_row_images)
                _shape_cell_text(
                    cell,
                    _shorten_words(val, limit),
                    body_size,
                    bold=(ci == 0),
                    align=PP_ALIGN.CENTER,
                    color=C_TEXT_DARK,
                    min_size=17,
                )
            x += w

    # Images are placed inside the Image column cell areas only when the Image column exists.
    if has_row_images and row_image_paths:
        image_col = n_cols - 1
        image_x = left + sum(col_widths[:image_col])
        col_w = col_widths[image_col]
        for ri, img_path in enumerate(row_image_paths[:len(rows)]):
            if not img_path or not os.path.exists(img_path):
                continue
            y = top + row_h * (ri + 1)
            pad = Inches(0.08)
            _add_image_contain(slide, img_path, image_x + pad, y + pad, col_w - pad * 2, row_h - pad * 2)

    return None



def create_presentation(logo_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def build_cover_slide(prs, chapter_name, chapter_number, logo_path, cover_image_path=None, slide_number=None, cover_image_paths=None, chapter_description=None):
    slide = _base_slide(prs)
    _white_bg(slide)
    _logo(slide, logo_path)
    _copyright(slide)
    _slide_number(slide, slide_number)

    title = _safe_slide_title(chapter_name, "")
    title_font = _fit_cover_font(title)
    _tb(slide, LEFT, Inches(1.20), COVER_TITLE_MAX_W, Inches(2.7),
        title, title_font, bold=True, color=C_TEXT_DARK, wrap=False)
    desc = _clean_text(chapter_description or "")
    if desc:
        _tb(slide, LEFT, Inches(4.25), Inches(8.9), Inches(0.9),
            _shorten_words(desc, 26), 26, italic=True, color=DARK_GRAY, wrap=True)
    _tb(slide, LEFT, Inches(5.60), Inches(8.8), Inches(0.45),
        f"Chapter {chapter_number}", 24, color=DARK_GRAY)
    _tb(slide, LEFT, Inches(8.55), Inches(6.0), Inches(0.45),
        "Lecture Slides", 24, color=DARK_GRAY)

    # Project override: first chapter/cover slide must use exactly ONE image.
    # No stacked image set and no repeated image.
    selected_cover_image = None
    if cover_image_path and os.path.exists(cover_image_path):
        selected_cover_image = cover_image_path
    elif cover_image_paths:
        for img in cover_image_paths:
            if img and os.path.exists(img):
                selected_cover_image = img
                break

    if selected_cover_image:
        _add_image_contain(slide, selected_cover_image, RIGHT_X, Inches(2.25), RIGHT_W, COVER_SINGLE_IMAGE_H)

    _notes(slide, "Welcome students. Introduce the lesson topic and outline the key concepts.")



def build_concept_slide(prs, lesson_name, body_text, sub_label=None,
                        image_path=None, speaker_notes=None, logo_path="", slide_number=None,
                        figure_legend=None, transition_caption=None, allow_no_image=False):
    slide = _base_slide(prs)
    _white_bg(slide)
    _logo(slide, logo_path)
    _copyright(slide)
    _slide_number(slide, slide_number)

    title = _safe_slide_title(lesson_name, body_text)
    if title.lower().startswith("definition") or title.lower() in {"core idea", "definition and core process"}:
        title = "What is the concept?"
    _tb(slide, LEFT, Inches(0.75), TEXT_W, Inches(1.15),
        title, FS_SLIDE_TITLE, bold=True, color=C_TEXT_DARK)
    _body_text(slide, body_text, Inches(2.55), max_words=65)
    if image_path:
        _image(slide, image_path, figure_legend=figure_legend)
    elif not allow_no_image:
        _image(slide, image_path)
    _transition_caption(slide, transition_caption)
    _notes(slide, speaker_notes)


def build_table_slide(prs, lesson_name, headers, rows, sub_title=None,
                      table_kind=None, image_path=None, row_image_paths=None, speaker_notes=None, logo_path="", slide_number=None,
                      figure_legend=None):
    slide = _base_slide(prs)
    _white_bg(slide)
    _logo(slide, logo_path)
    _copyright(slide)
    _slide_number(slide, slide_number)

    seed_text = " ".join(" ".join(map(str, r)) if isinstance(r, list) else str(r) for r in (rows or [])[:2])
    title = _safe_slide_title(sub_title or lesson_name, seed_text)
    _tb(slide, LEFT, Inches(0.42), Inches(17.25), TITLE_BOX_H,
        title, FS_SLIDE_TITLE, bold=True, color=C_TEXT_DARK, wrap=True)

    _add_table(slide, headers, rows, Inches(0.75), TITLE_SAFE_TOP_TABLE, Inches(18.5),
               max_height=Inches(7.8), row_image_paths=row_image_paths, max_rows=4)
    _notes(slide, speaker_notes)


def _discussion_header(slide):
    # Exact reference layout measured from the approved Natural Selection discussion template.
    _tb(slide, Inches(1.0417), Inches(0.8765), Inches(18.4541), Inches(0.7415),
        "Discussion", FS_SLIDE_TITLE, bold=True, color=C_TEXT_DARK)

    icon_path = _discussion_icon_path()
    if icon_path:
        slide.shapes.add_picture(icon_path, Inches(1.0417), Inches(3.3282), width=Inches(0.3750), height=Inches(0.3750))

    _tb(slide, Inches(1.1947), Inches(2.9993), Inches(5.0862), Inches(1.0000),
        "Discuss with the class", 30, color=C_ACCENT_BLUE)



def build_discussion_question_slide(prs, lesson_name, question_text,
                                     hint_text=None, image_path=None,
                                     speaker_notes=None, logo_path="", slide_number=None,
                                     figure_legend=None):
    slide = _base_slide(prs)
    _white_bg(slide)
    _logo(slide, logo_path)
    _copyright(slide)
    _slide_number(slide, slide_number)

    _discussion_header(slide)

    # Fixed non-overlap layout with complete-sentence visible text.
    # Avoid false sentence fragments caused by blunt word truncation.
    question_clean = _complete_sentence_text(question_text, "", max_words=34, prefer_question=True)
    hint_clean = _complete_sentence_text(hint_text, "", max_words=24) if hint_text else ""

    _tb(slide, Inches(1.0417), Inches(3.70), Inches(9.35), Inches(2.55),
        question_clean, 36, bold=True, color=C_TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True)

    if hint_clean:
        _tb(slide, Inches(1.0417), Inches(6.55), Inches(9.35), Inches(1.25),
            "Hint: " + hint_clean, FS_BODY_SECONDARY, italic=True, color=DARK_GRAY, align=PP_ALIGN.LEFT, wrap=True)

    _image(slide, image_path, figure_legend=figure_legend)
    _notes(slide, speaker_notes)



def build_discussion_answer_slide(prs, lesson_name, answer_summary,
                                   answer_explanation, image_path=None,
                                   speaker_notes=None, logo_path="", slide_number=None,
                                   figure_legend=None):
    slide = _base_slide(prs)
    _white_bg(slide)
    _logo(slide, logo_path)
    _copyright(slide)
    _slide_number(slide, slide_number)

    _discussion_header(slide)
    answer_headline = _complete_sentence_text(answer_summary, answer_explanation, max_words=18)
    answer_body = _complete_sentence_text(answer_explanation, answer_summary, max_words=44)
    _tb(slide, Inches(1.0417), Inches(3.88), Inches(9.35), Inches(1.65),
        "Answer: " + answer_headline, 31, bold=True, color=C_TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True)
    _tb(slide, Inches(0.9561), Inches(5.80), Inches(9.45), Inches(2.75),
        answer_body, 28, color=C_TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True)
    _image(slide, image_path, figure_legend=figure_legend)
    _notes(slide, speaker_notes)



def build_summary_slide(prs, summary_statement, table_headers=None,
                         table_rows=None, logo_path="", speaker_notes=None, slide_number=None):
    slide = _base_slide(prs)
    _white_bg(slide)
    _logo(slide, logo_path)
    _copyright(slide)
    _slide_number(slide, slide_number)

    _tb(slide, LEFT, Inches(0.55), Inches(17.0), Inches(0.95),
        "Summary", FS_SLIDE_TITLE, bold=True, color=C_TEXT_DARK)

    clean_summary = _complete_sentence_text(_clean_text(summary_statement), "", max_words=26)
    _tb(slide, LEFT, Inches(1.45), Inches(17.0), Inches(0.72),
        clean_summary, 26, italic=True, color=C_TEXT_DARK)

    rows = (table_rows or [])[:5]
    if table_headers and rows:
        # Same visual-grid table, but with more vertical room to avoid text overlap.
        _add_table(slide, table_headers, rows, LEFT, Inches(2.30), Inches(17.0),
                   max_height=Inches(6.85), row_image_paths=None, max_rows=5)
    _notes(slide, speaker_notes)



def build_glossary_slide(prs, terms_dict, logo_path="", slide_number=None):
    slide = _base_slide(prs)
    _white_bg(slide)
    _logo(slide, logo_path)
    _copyright(slide)
    _slide_number(slide, slide_number)

    _tb(slide, LEFT, Inches(0.65), Inches(17.0), Inches(0.9),
        "Glossary", FS_SLIDE_TITLE, bold=True, color=C_TEXT_DARK)

    items = list(terms_dict.items())[:9]
    box = slide.shapes.add_textbox(LEFT, Inches(1.70), Inches(17.2), Inches(8.2))
    tf = box.text_frame
    tf.clear()
    _lock_text_frame(tf, vertical_anchor=MSO_ANCHOR.TOP, wrap=True)
    first = True
    for term, definition in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _format_paragraph(p, align=PP_ALIGN.LEFT, space_before=(6 if p is not tf.paragraphs[0] else 0), space_after=0, line_spacing=1.0)
        r1 = p.add_run()
        r1.text = f"{_clean_text(term)}: "
        _font(r1, FS_BODY_SECONDARY, bold=True, color=C_TEXT_DARK)
        r2 = p.add_run()
        r2.text = _complete_sentence_text(definition, "", max_words=16)
        _font(r2, FS_BODY_SECONDARY, color=C_TEXT_DARK)
    _notes(slide, "Review these key terms with students.")
